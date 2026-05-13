"""
PowerPoint Generator: Evaluation Paradigms in Recommender Systems
B.Tech CSE - 6th Semester - Recommender Systems

Run this script: python create_ppt.py
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # Blank slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Add background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(0, 51, 102)  # Dark blue
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RgbColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_list, is_two_column=False):
    slide_layout = prs.slide_layouts[6]  # Blank slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RgbColor(0, 51, 102)
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    
    if is_two_column and len(content_list) == 2:
        # Two column layout
        for i, (col_title, items) in enumerate(content_list):
            x_pos = Inches(0.5) if i == 0 else Inches(5)
            content_box = slide.shapes.add_textbox(x_pos, Inches(1.5), Inches(4.5), Inches(5))
            tf = content_box.text_frame
            tf.word_wrap = True
            
            # Column title
            p = tf.paragraphs[0]
            p.text = col_title
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RgbColor(0, 51, 102)
            
            for item in items:
                p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(16)
                p.font.color.rgb = RgbColor(50, 50, 50)
                p.space_before = Pt(8)
    else:
        # Single column
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(content_list):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(20)
            p.font.color.rgb = RgbColor(50, 50, 50)
            p.space_before = Pt(12)

def add_formula_slide(prs, title, formulas):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RgbColor(0, 51, 102)
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    
    y_pos = 1.6
    for name, formula, description in formulas:
        # Formula name
        name_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.4))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RgbColor(0, 102, 153)
        
        # Formula box
        formula_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y_pos + 0.4), Inches(8.4), Inches(0.5))
        formula_shape.fill.solid()
        formula_shape.fill.fore_color.rgb = RgbColor(240, 240, 240)
        formula_shape.line.color.rgb = RgbColor(200, 200, 200)
        
        formula_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos + 0.45), Inches(8), Inches(0.4))
        tf = formula_box.text_frame
        p = tf.paragraphs[0]
        p.text = formula
        p.font.size = Pt(14)
        p.font.name = "Consolas"
        p.font.color.rgb = RgbColor(50, 50, 50)
        p.alignment = PP_ALIGN.CENTER
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos + 0.95), Inches(8.4), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = description
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = RgbColor(100, 100, 100)
        
        y_pos += 1.5

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs, 
        "Evaluation Paradigms in\nRecommender Systems",
        "B.Tech CSE - 6th Semester\nRecommender Systems")
    
    # Slide 2: Table of Contents
    add_content_slide(prs, "Agenda", [
        "Introduction to Recommender System Evaluation",
        "Offline Evaluation Methods",
        "Online Evaluation Methods",
        "User Studies & Qualitative Evaluation",
        "Accuracy Metrics (Rating & Ranking)",
        "Beyond Accuracy Metrics",
        "Challenges & Best Practices",
        "Conclusion"
    ])
    
    # Slide 3: Introduction
    add_content_slide(prs, "Why Evaluate Recommender Systems?", [
        "Measure how well recommendations match user preferences",
        "Compare different algorithms and approaches",
        "Identify areas for improvement",
        "Ensure business goals are met (engagement, revenue)",
        "Balance multiple objectives (accuracy, diversity, fairness)",
        "Validate system before deployment"
    ])
    
    # Slide 4: Three Evaluation Paradigms
    add_content_slide(prs, "Three Evaluation Paradigms", [
        ("Offline Evaluation", [
            "Uses historical data",
            "Fast and reproducible",
            "No user involvement needed",
            "May not reflect real behavior"
        ]),
        ("Online Evaluation", [
            "Real-time user interactions",
            "A/B testing & live experiments",
            "Directly measures impact",
            "Expensive and time-consuming"
        ])
    ], is_two_column=True)
    
    # Slide 5: User Studies
    add_content_slide(prs, "User Studies & Qualitative Evaluation", [
        "Controlled experiments with real users",
        "Surveys and questionnaires for user satisfaction",
        "Think-aloud protocols to understand user reasoning",
        "Focus groups for in-depth feedback",
        "Measures subjective qualities: trust, satisfaction, novelty",
        "Expensive but provides deep insights into user experience"
    ])
    
    # Slide 6: Offline Evaluation Process
    add_content_slide(prs, "Offline Evaluation Process", [
        "Data Splitting: Train/Validation/Test sets",
        "Cross-Validation: k-fold for robust estimates",
        "Temporal Split: Train on past, test on future",
        "Leave-One-Out: Remove one interaction per user",
        "Holdout Methods: Random vs. temporal sampling",
        "Important: Avoid data leakage from future to past!"
    ])
    
    # Slide 7: Rating Prediction Metrics
    add_formula_slide(prs, "Rating Prediction Metrics", [
        ("Mean Absolute Error (MAE)", 
         "MAE = (1/n) × Σ|r_ui - r̂_ui|",
         "Average absolute difference between predicted and actual ratings"),
        ("Root Mean Square Error (RMSE)", 
         "RMSE = √[(1/n) × Σ(r_ui - r̂_ui)²]",
         "Penalizes larger errors more heavily than MAE"),
        ("Normalized MAE", 
         "NMAE = MAE / (r_max - r_min)",
         "Scale-independent metric for comparing across datasets")
    ])
    
    # Slide 8: Ranking Metrics
    add_formula_slide(prs, "Ranking Metrics", [
        ("Precision@K", 
         "P@K = |Relevant ∩ Recommended@K| / K",
         "Fraction of recommended items that are relevant"),
        ("Recall@K", 
         "R@K = |Relevant ∩ Recommended@K| / |Relevant|",
         "Fraction of relevant items that are recommended"),
        ("F1-Score@K", 
         "F1@K = 2 × (P@K × R@K) / (P@K + R@K)",
         "Harmonic mean of precision and recall")
    ])
    
    # Slide 9: Advanced Ranking Metrics
    add_formula_slide(prs, "Advanced Ranking Metrics", [
        ("Mean Average Precision (MAP)", 
         "MAP = (1/|U|) × Σ AP(u)",
         "Average of Average Precision across all users"),
        ("Normalized DCG (NDCG)", 
         "NDCG@K = DCG@K / IDCG@K",
         "Measures ranking quality with position-based discounting"),
        ("Mean Reciprocal Rank (MRR)", 
         "MRR = (1/|U|) × Σ(1/rank_i)",
         "Average reciprocal position of first relevant item")
    ])
    
    # Slide 10: Beyond Accuracy - Coverage & Diversity
    add_content_slide(prs, "Beyond Accuracy: Coverage & Diversity", [
        ("Coverage Metrics", [
            "Catalog Coverage: % of items recommended",
            "User Coverage: % of users receiving recommendations",
            "Prediction Coverage: % of items with predictions"
        ]),
        ("Diversity Metrics", [
            "Intra-List Diversity: Variety within a list",
            "Inter-List Diversity: Variety across users",
            "Aggregate Diversity: Overall system diversity"
        ])
    ], is_two_column=True)
    
    # Slide 11: Beyond Accuracy - Novelty & Serendipity
    add_content_slide(prs, "Beyond Accuracy: Novelty & Serendipity", [
        "Novelty: Recommending items users haven't seen before",
        "Serendipity: Surprising and useful recommendations",
        "Long-tail items: Less popular but potentially interesting",
        "Trade-off: Popular items are safer but less novel",
        "Serendipity = Relevance × Unexpectedness",
        "Important for user satisfaction and discovery"
    ])
    
    # Slide 12: Fairness in Evaluation
    add_content_slide(prs, "Fairness & Bias in Evaluation", [
        "User Fairness: Equal quality across user groups",
        "Item Fairness: Fair exposure for all items/providers",
        "Popularity Bias: Over-recommending popular items",
        "Position Bias: Users click top results regardless of relevance",
        "Selection Bias: Training data reflects past recommendations",
        "Fairness-Accuracy Trade-offs must be considered"
    ])
    
    # Slide 13: Online Evaluation - A/B Testing
    add_content_slide(prs, "Online Evaluation: A/B Testing", [
        "Randomly assign users to control (A) or treatment (B) group",
        "Control group: Current system",
        "Treatment group: New recommendation algorithm",
        "Measure key metrics: CTR, conversion, engagement, revenue",
        "Statistical significance testing (p-values, confidence intervals)",
        "Run for sufficient time to capture variability"
    ])
    
    # Slide 14: Online Metrics
    add_content_slide(prs, "Online Evaluation Metrics", [
        "Click-Through Rate (CTR): Clicks / Impressions",
        "Conversion Rate: Purchases / Clicks",
        "Dwell Time: Time spent on recommended items",
        "Session Length: Duration of user engagement",
        "Return Rate: Users coming back",
        "Revenue per User: Direct business impact"
    ])
    
    # Slide 15: Challenges
    add_content_slide(prs, "Challenges in Evaluation", [
        "Cold Start: New users/items lack interaction data",
        "Sparsity: Most user-item pairs are unobserved",
        "Feedback Loop: Evaluating on biased historical data",
        "Offline-Online Gap: Offline metrics may not predict online success",
        "Multi-objective Trade-offs: Accuracy vs. diversity vs. fairness",
        "Scalability: Evaluating on large-scale systems"
    ])
    
    # Slide 16: Best Practices
    add_content_slide(prs, "Best Practices for Evaluation", [
        "Use multiple metrics - don't rely on just one",
        "Consider both accuracy and beyond-accuracy metrics",
        "Use appropriate data splitting (temporal for time-sensitive)",
        "Report statistical significance and confidence intervals",
        "Combine offline evaluation with online A/B tests",
        "Document evaluation protocol for reproducibility"
    ])
    
    # Slide 17: Summary
    add_content_slide(prs, "Summary", [
        "Evaluation is crucial for building effective recommender systems",
        "Three paradigms: Offline, Online, and User Studies",
        "Accuracy metrics: MAE, RMSE, Precision, Recall, NDCG",
        "Beyond accuracy: Coverage, Diversity, Novelty, Fairness",
        "No single metric captures everything - use multiple",
        "Offline evaluation is fast; online evaluation is ground truth"
    ])
    
    # Slide 18: References
    add_content_slide(prs, "References", [
        "Ricci, Rokach, Shapira - Recommender Systems Handbook",
        "Herlocker et al. - Evaluating Collaborative Filtering",
        "Shani & Gunawardana - Evaluating Recommendation Systems",
        "Jannach et al. - Measuring the Business Value of RS",
        "Burke et al. - Fairness in Recommendation",
        "Cremonesi et al. - Performance of Recommender Algorithms"
    ])
    
    # Slide 19: Thank You
    add_title_slide(prs,
        "Thank You!",
        "Questions?")
    
    # Save presentation
    output_path = "Evaluation_Paradigms_Recommender_Systems.pptx"
    prs.save(output_path)
    print(f"\n✅ Presentation saved as: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")
    return output_path

if __name__ == "__main__":
    print("Creating PowerPoint Presentation...")
    print("Subject: Recommender Systems")
    print("Topic: Evaluation Paradigms in Recommender Systems")
    print("-" * 50)
    create_presentation()
